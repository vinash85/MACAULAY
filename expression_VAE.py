import torch
import torch.nn as nn
import torch.optim as optim
from sklearn.preprocessing import StandardScaler
import pandas as pd
from torch.utils.data import DataLoader, TensorDataset
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

prs = Presentation()
#import the transposed omics expression data
df = pd.read_csv('./datas/omic_expression_transpose.csv.csv')
# Extract the numeric values from the DataFrame
numeric_data = df.iloc[1:, 1:].values
scaler = StandardScaler()
scaled_data = scaler.fit_transform(numeric_data)

tensor_data = torch.tensor(scaled_data, dtype=torch.float32)
#print(tensor_data)
print('tensor data shape:', tensor_data.shape)
print("df shape:", df.shape)
class VAE(nn.Module):
    def __init__(self, input_dim, latent_dim, hidden_dim):
        super(VAE, self).__init__()

        self.input_dim = input_dim
        self.latent_dim = latent_dim
        self.hidden_dim = hidden_dim
        #Encoded layer
        self.encoder = nn.Sequential(
            nn.Linear(input_dim, hidden_dim),
            nn.ReLU()
        )

        # Latent space layers
        self.fc_mu = nn.Linear(hidden_dim, latent_dim)
        self.fc_logvar = nn.Linear(hidden_dim, latent_dim)

        # Decoder layers
        self.decoder = nn.Sequential(
            nn.Linear(latent_dim, hidden_dim),
            nn.ReLU(),
            nn.Linear(hidden_dim, input_dim),
            nn.Sigmoid()  
        )

    def encode(self, x):
        hidden = self.encoder(x)
        mu = self.fc_mu(hidden)
        logvar = self.fc_logvar(hidden)
        return mu, logvar

    def reparameterize(self, mu, logvar):
        std = torch.exp(0.5 * logvar)
        eps = torch.randn_like(std)
        z = mu + eps * std
        return z

    def decode(self, z):
        output = self.decoder(z)
        return output

    def forward(self, x):
        mu, logvar = self.encode(x)
        z = self.reparameterize(mu, logvar)
        output = self.decode(z)
        return output, mu, logvar

device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
tensor_data = tensor_data.to(device)

# Define hyperparameters
input_dim = tensor_data.shape[1]  # Dimensionality of the omic gene expression data
learning_rate = [0.0001, 0.001]
batch_size = [16, 32, 64, 128]
num_epochs = 50

for hidden_i in range(200, 1000, 100):

    for l_i, learning_i in enumerate(learning_rate):

        for b_i, batch_i in enumerate(batch_size):

            # Create an instance of the VAE
            latent_dim = hidden_i - 100
            vae = VAE(input_dim, latent_dim, hidden_i).to(device)

            # Define reconstruction loss (e.g., Mean Squared Error)
            reconstruction_loss = nn.MSELoss()

            # Define optimizer
            optimizer = optim.Adam(vae.parameters(), lr=learning_i)

            # Create DataLoader for batch training
            dataset = TensorDataset(tensor_data)
            dataloader = DataLoader(dataset, batch_size=batch_i, shuffle=True)

            # Training loop
            for epoch in range(num_epochs):
                running_loss = 0.0

                for batch_data in dataloader:
                    # Zero the gradients
                    optimizer.zero_grad()

                    # Move batch data to the appropriate device
                    batch_data = batch_data[0].to(device)

                    # Forward pass
                    output, mu, logvar = vae(batch_data)

                    # Compute reconstruction loss
                    recon_loss = reconstruction_loss(output, batch_data)

                    # Compute KL divergence loss
                    kl_loss = -0.5 * torch.sum(1 + logvar - mu.pow(2) - logvar.exp())

                    # Total loss
                    total_loss = recon_loss + kl_loss

                    # Backward pass and optimization
                    total_loss.backward()
                    optimizer.step()

                    # Accumulate running loss
                    running_loss += total_loss.item()

                # Print epoch statistics
                epoch_loss = running_loss / len(dataloader)
                print(f"Epoch [{epoch + 1}/{num_epochs}], Loss: {epoch_loss:.4f}, hidden_dim: {hidden_i}, learning_rate: {learning_i}, batch_size: {batch_i}")
                with open('epoch.txt', 'a') as f:
                    f.write(f"Epoch [{epoch + 1}/{num_epochs}], Loss: {epoch_loss:.4f}, hidden_dim: {hidden_i}, learning_rate: {learning_i}, batch_size: {batch_i}")
                    f.write('\n')

                if epoch == num_epochs-1:
                    with open('last_epoch.txt', 'a') as f:
                        f.write(f"Epoch [{epoch + 1}/{num_epochs}], Loss: {epoch_loss:.4f}, hidden_dim: {hidden_i}, learning_rate: {learning_i}, batch_size: {batch_i}")
                        f.write('\n')
                    # torch.save(vae.state_dict(), f'./models/{hidden_i}_{learning_i}_{batch_i}.pth')
            # Reconstruction example
            sample_data = tensor_data[:10]  # Taking the first 10 samples as an example
            sample_data = sample_data.to(device)
            reconstructed_data, _, _ = vae(sample_data)
            # print("Original Data:")
            # print(sample_data)
            # print("Reconstructed Data:")
            # print(reconstructed_data)
            

            # Convert the tensors back to numpy arrays for plotting
            sample_data_np = sample_data.cpu().detach().numpy()
            reconstructed_data_np = reconstructed_data.cpu().detach().numpy()

            # Plotting the original and reconstructed data
            fig, axes = plt.subplots(nrows=2, ncols=1, figsize=(10, 6))

            # Plot original data
            axes[0].imshow(sample_data_np, aspect='auto', cmap='viridis')
            axes[0].set_title('Original Data')

            # Plot reconstructed data
            axes[1].imshow(reconstructed_data_np, aspect='auto', cmap='viridis')
            axes[1].set_title('Reconstructed Data')

            # Adjust the spacing between subplots
            plt.tight_layout()

            # Show the plot
            plt.title(f'hidden_dim: {hidden_i}, learning_rate: {learning_i}, batch_size: {batch_i}')
            plt.savefig(f'plot_images/hidden_{hidden_i}_learning_{learning_i}_batch_{batch_i}_plot.png')

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(f'plot_images/hidden_{hidden_i}_learning_{learning_i}_batch_{batch_i}_plot.png' , Inches(1), Inches(1), width=Inches(8), height=Inches(6))

prs.save(f'slides/Deep_Plot_Essentiality.pptx')
            




