import torch
import torch.nn as nn
import torch.optim as optim
from sklearn.preprocessing import StandardScaler
from sklearn.model_selection import train_test_split
import pandas as pd
from torch.utils.data import DataLoader, TensorDataset
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np

prs = Presentation()
#import the transposed omics expression data
df = pd.read_csv('OmicsExpressionProteinCodingGenesTPMLogp1.csv')
# Extract the numeric values from the DataFrame
numeric_data = df.iloc[1:, 1:].values
scaler = StandardScaler()
scaled_data = scaler.fit_transform(numeric_data)

# Split the data into train and test sets
train_data, test_data = train_test_split(scaled_data, test_size=0.2, random_state=42)

train_tensor = torch.tensor(train_data, dtype=torch.float32)
test_tensor = torch.tensor(test_data, dtype=torch.float32)


#print(train_tensor)
print("df shape:", df.shape)
print('training tensor data shape:', train_tensor.shape)

class VAE(nn.Module):
    def __init__(self, input_dim, latent_dim, hidden_dim):
        super(VAE, self).__init__()

        self.input_dim = input_dim
        self.latent_dim = latent_dim
        self.hidden_dim = hidden_dim
        #Encoded layer
        self.encoder = nn.Sequential(
            nn.Linear(input_dim, hidden_dim),
            nn.ReLU()
        )

        # Latent space layers
        self.fc_mu = nn.Linear(hidden_dim, latent_dim)
        self.fc_logvar = nn.Linear(hidden_dim, latent_dim)

        # Decoder layers
        self.decoder = nn.Sequential(
            nn.Linear(latent_dim, hidden_dim),
            nn.ReLU(),
            nn.Linear(hidden_dim, input_dim),
            nn.Sigmoid()  
        )

    def encode(self, x):
        hidden = self.encoder(x)
        mu = self.fc_mu(hidden)
        logvar = self.fc_logvar(hidden)
        return mu, logvar

    def reparameterize(self, mu, logvar):
        std = torch.exp(0.5 * logvar)
        eps = torch.randn_like(std)
        z = mu + eps * std
        return z

    def decode(self, z):
        output = self.decoder(z)
        return output

    def forward(self, x):
        mu, logvar = self.encode(x)
        z = self.reparameterize(mu, logvar)
        output = self.decode(z)
        return output, mu, logvar

device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
train_tensor = train_tensor.to(device)

# Define hyperparameters
input_dim = train_tensor.shape[1]  # Dimensionality of the omic gene expression data
learning_rate = [0.0001]
batch_size = [128]
num_epochs = 2000

for hidden_i in range(600, 700, 100):

    for l_i, learning_i in enumerate(learning_rate):

        for b_i, batch_i in enumerate(batch_size):

            # Create an instance of the VAE
            latent_dim = hidden_i - 100
            vae = VAE(input_dim, latent_dim, hidden_i).to(device)

            # Define reconstruction loss (e.g., Mean Squared Error)
            reconstruction_loss = nn.MSELoss()

            # Define optimizer
            optimizer = optim.Adam(vae.parameters(), lr=learning_i)

            # Create DataLoader for batch training
            dataset = TensorDataset(train_tensor)
            dataloader = DataLoader(dataset, batch_size=batch_i, shuffle=True)

            # Training loop
            for epoch in range(num_epochs):
                running_recon_loss = 0.0


                for batch_data in dataloader:
                    # Zero the gradients
                    optimizer.zero_grad()

                    # Move batch data to the appropriate device
                    batch_data = batch_data[0].to(device)

                    # Forward pass
                    output, mu, logvar = vae(batch_data)

                    # Compute reconstruction loss
                    recon_loss = reconstruction_loss(output, batch_data)

                    # Backward pass and optimization
                    recon_loss.backward()
                    optimizer.step()

                    # Accumulate running reconstruction loss
                    running_recon_loss += recon_loss.item()

                # Print epoch statistics
                epoch_recon_loss = running_recon_loss / len(dataloader)
                print(f"Epoch [{epoch + 1}/{num_epochs}], Training Reconstruction Loss: {epoch_recon_loss:.4f}, hidden_dim: {hidden_i}, learning_rate: {learning_i}, batch_size: {batch_i}")

                if epoch % 200 == 0:
                    # Save the trained model
                    # torch.save(vae.state_dict(), f"vae_state/vae_{epoch}.pth")
                    # Pass the test data through the trained VAE model
                    reconstructed_test_data, _, _ = vae(test_tensor)

                    # Convert the tensors back to numpy arrays
                    test_data_np = test_tensor.cpu().detach().numpy()
                    reconstructed_test_data_np = reconstructed_test_data.cpu().detach().numpy()

                    # Calculate the reconstruction loss for the test data
                    test_recon_loss = reconstruction_loss(reconstructed_test_data, test_tensor)

                    # Print the reconstruction loss for the test data
                    print(f"Reconstruction Loss (Test Data) {epoch}: {test_recon_loss:.4f}")

                    # Calculate the correlation coefficients for all test genes
                    correlation_coefficients = []
                    for i in range(test_tensor.shape[1]):
                        test_gene = test_tensor[:, i].cpu().detach().numpy()
                        reconstructed_test_gene = reconstructed_test_data[:, i].cpu().detach().numpy()
                        correlation = np.corrcoef(test_gene, reconstructed_test_gene)[0, 1]
                        correlation_coefficients.append(correlation)

                    # Plot the correlation graph
                    plt.figure()
                    plt.scatter(range(test_tensor.shape[1]), correlation_coefficients)
                    plt.xlabel('Gene Index')
                    plt.ylabel('Correlation Coefficient')
                    plt.title(f'Correlation of Reconstructed Genes at epoch {epoch}')
                    # plt.show()
                    
                    plt.savefig(f'plot_images/Correlation of Reconstructed Genes (Test Data) {epoch}.png')
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(f'plot_images/Correlation of Reconstructed Genes (Test Data) {epoch}.png' , Inches(1), Inches(1), width=Inches(8), height=Inches(6))


                    # Plot the histogram of correlation coefficients
                    plt.figure()
                    plt.hist(correlation_coefficients, bins=30, edgecolor='black')
                    plt.xlabel('Correlation Coefficient')
                    plt.ylabel('Frequency')
                    plt.title(f'Histogram of Correlation Coefficients at epoch {epoch}')
                    # plt.show()
                    plt.savefig(f'plot_images/Histogram of Correlation Coefficients{epoch}.png')

                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(f'plot_images/Histogram of Correlation Coefficients{epoch}.png' , Inches(1), Inches(1), width=Inches(8), height=Inches(6))

                    
            








            # Reconstruction example
            # Pass the test data through the trained VAE model
            torch.save(vae.state_dict(), f"vae_state/vae_expression.pth")
            reconstructed_test_data, _, _ = vae(test_tensor)

            # Convert the tensors back to numpy arrays
            test_data_np = test_tensor.cpu().detach().numpy()
            reconstructed_test_data_np = reconstructed_test_data.cpu().detach().numpy()

            # Calculate the reconstruction loss for the test data
            test_recon_loss = reconstruction_loss(reconstructed_test_data, test_tensor)

            # Print the reconstruction loss for the test data
            print(f"Reconstruction Loss (Test Data): {test_recon_loss:.4f}")

            # Calculate the correlation coefficients for all test genes
            correlation_coefficients = []
            for i in range(test_tensor.shape[1]):
                test_gene = test_tensor[:, i].cpu().detach().numpy()
                reconstructed_test_gene = reconstructed_test_data[:, i].cpu().detach().numpy()
                correlation = np.corrcoef(test_gene, reconstructed_test_gene)[0, 1]
                correlation_coefficients.append(correlation)

            # Plot the correlation graph
            plt.figure()
            plt.scatter(range(test_tensor.shape[1]), correlation_coefficients)
            plt.xlabel('Gene Index')
            plt.ylabel('Correlation Coefficient')
            plt.title(f'Correlation of Reconstructed Genes (Test Data)')
            # plt.show()
            
            plt.savefig(f'figures/Correlation of Reconstructed Genes (Test Data)')

            # Plot the histogram of correlation coefficients
            plt.figure()
            plt.hist(correlation_coefficients, bins=30, edgecolor='black')
            plt.xlabel('Correlation Coefficient')
            plt.ylabel('Frequency')
            plt.title(f'Histogram of Correlation Coefficients')
            # plt.show()
            plt.savefig(f'figures/Histogram of Correlation Coefficients')

prs.save(f'slide.pptx')
            









